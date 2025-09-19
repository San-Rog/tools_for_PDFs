import pymupdf
import streamlit as st
import streamlit.components.v1 as components
import zipfile
import os
import re
import io
import time
import unidecode
import textwrap
import xlsxwriter
import numpy as np
import pandas as pd
import random
from segno import helpers
import subprocess
import datetime
from PyPDF2 import PdfReader, PdfWriter
from pdf_watermark_remover import process_pdf
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from brutils import is_valid_email
from brutils import is_valid_phone
from brutils import remove_symbols_phone
import itertools
    
@st.cache_data   
def nameFile():
    symbols = ['-', ':', '.']
    nowTime = str(datetime.datetime.now())
    try:
        for symbol in symbols: 
            nowTime = nowTime.replace(symbol, '_')
    except:
        pass
    return nowTime
    
@st.cache_data  
def extractText(filePdf, mode):
    text = ''
    docPdf = pymupdf.open(filePdf)
    for page in docPdf:
        if mode == 0:
            text += page.get_text('text')
        elif mode == 1:
            text += page.get_text('html')
        else:
            text += page.get_text('xhtml')
    docPdf.close()
    return text
        
@st.cache_data
def extractUrls(filePdf):
    docPdf = pymupdf.open(filePdf)
    allLinks = []
    for p, page in enumerate(docPdf):
        links = page.get_links()
        for link in links:
            try:
                nameUrl = link["uri"]
                fromUrl = link["from"]
                newText = f'{nameUrl}; coordenadas: {fromUrl}\n'
                allLinks.append(newText)
            except:
                pass
    text = ''.join(allLinks) 
    docPdf.close()
    return text
    
def mensResult(value, nFiles, modelButt, fileTmp, fileFinal):
    opt = st.session_state[listKeys[5]]
    if opt == 0:
        crt = optionsSel[3]
    else:
        if opt != - 1:
            mult = st.session_state[listKeys[6]]
            if mult == 0:
                crt = f'{optionsSel[opt]}' 
            else:
                crt = f'{optionsSel[opt]} {st.session_state[listKeys[6]]}' 
        else: 
          crt = f'segmentação com base em {st.session_state[listKeys[6]]} página(s)'  
    colIcoLim, colMens, colDown = st.columns([1, 16, 4], vertical_alignment='bottom', 
                                                         width='stretch')
    textIcoLim = colIcoLim.text_input(icons[6][0], icon=icons[6][1], width=40)
    if value == 1:
        if modelButt == 'zip': 
                with open(fileTmp, "rb") as file:
                    colDown.download_button(label="Download",
                                            data=file,
                                            file_name=fileFinal,
                                            mime='application/zip', 
                                            icon=":material/download:", 
                                            use_container_width=True)
        colMens.success(f'Gerado o zipado :blue[**{fileFinal}**] com ***{nFiles}*** arquivo(s) (:red[**{crt}**]). Clique no botão ao lado 👉.', 
                        icon='✔️') 
    elif value == 0:
        colDown.download_button(label='Download', 
                                data=fileTmp,
                                file_name=fileFinal,
                                mime='application/octet-stream', 
                                icon=":material/download:", 
                                use_container_width=True)
        colMens.success(f'Gerado o arquivo :blue[**{fileFinal}**] (:red[**{crt}**]). Clique no botão ao lado 👉.', 
                        icon='✔️') 
    elif value == 2:
        colDown.download_button(label='Download',
                                data=fileTmp,
                                file_name=fileFinal,
                                mime="text/csv", 
                                icon=":material/download:", 
                                use_container_width=True)
        colMens.success(f'Gerado o arquivo :blue[**{fileFinal}**] (:red[**{crt}**]). Clique no botão ao lado 👉.', 
                        icon='✔️')
    elif value == 3:
        colDown.download_button(label='Download',
                                data=fileTmp,
                                file_name=fileFinal,
                                mime='application/octet-stream', 
                                use_container_width=True)
        colMens.success(f'Gerado o arquivo :blue[**{fileFinal}**] (:red[**{crt}**]). Clique no botão ao lado 👉.', 
                        icon='✔️')
    elif value == 4:
        colDown.download_button(label='Download',
                                data=fileTmp,
                                file_name=fileFinal,
                                 mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document', 
                                use_container_width=True)
        colMens.success(f'Gerado o arquivo :blue[**{fileFinal}**] (:red[**{crt}**]). Clique no botão ao lado 👉.', 
                        icon='✔️')
    upDownScroll(6)
    
def extractImgs(filePdf):
    docPdf = pymupdf.open(filePdf)
    allImgName = []
    for p, page in enumerate(docPdf):
        imageList = page.get_images(full=True)
        if imageList:
            for i, img in enumerate(imageList):
                xref = img[0]
                baseImg = docPdf.extract_image(xref)
                imgBytes = baseImg["image"]
                imgExt = baseImg["ext"]
                imgName = f"image_{p+1}_{i+1}.{imgExt}"
                with open(imgName, "wb") as fileImg:
                    fileImg.write(imgBytes)
                allImgName.append(imgName)
    return allImgName

def downloadExt(files, namePdf, numPgOne, numPgTwo, obj):
    fileTmp = f'{nameFile()}_tempFile.zip'
    fileZip = f'{namePdf}_{numPgOne}_{numPgTwo}_{nameFile()}.zip'
    for file in files:
        with open(file, "rb") as extFile:
           PDFbyte = extFile.read()
        with zipfile.ZipFile(fileTmp, 'a') as extFile:
           extFile.writestr(file, PDFbyte)
    nFiles = len(files) 
    if nFiles > 0:
        mensResult(1, len(files), 'zip', fileTmp, fileZip)
    else:
        strEmpty = f"😢 Operação fracassada para '{obj}' do arquivo '{namePdf}', intervalo de páginas {numPgOne}-{numPgTwo}!"
        config(strEmpty)

def rotatePdf(filePdf, index):
    inputPdf = filePdf
    name, ext = os.path.splitext(inputPdf)
    angle = int(valAngles[index].replace('°', ''))
    output = f'{name}_rotate_{angle}{ext}'
    docPdf = pymupdf.open(filePdf)
    for page in docPdf:
        page.set_rotation(angle)
    docPdf.save(output)
    docPdf.close()
    return output
    
def saveAllPdf(outputBase, contPartes, writer):
    outputPdf = f"{outputBase}{contPartes + 1}.pdf"
    with open(outputPdf, "wb") as outputFile:
        writer.write(outputFile)
    docPdf = pymupdf.open(outputFile)
    countPg.append(len(docPdf))
    docPdf.close()
    return outputPdf

def divideBySize(inputPdf, sizeMax, outputBase):
    filesCutSave = []
    try:
        reader = PdfReader(inputPdf)
        nPgs = len(reader.pages)
        sizeActual = 0
        contPartes = 0
        writer = PdfWriter()
        for i in range(nPgs):
            nameTeste = f'teste_{i+1}.pdf'
            namesTeste.append(nameTeste)
            page = reader.pages[i]
            writer.add_page(page) 
            with open(nameTeste, 'wb') as g:
                writer.write(g)
            sizeActual = os.path.getsize(nameTeste)/(1024**2)
            if sizeActual >= sizeMax:
                outputPdf = saveAllPdf(outputBase, contPartes, writer)
                filesCutSave.append(outputPdf)
                writer = PdfWriter()
                sizeActual = 0
                contPartes += 1
        if len(writer.pages) > 0:
            outputPdf = saveAllPdf(outputBase, contPartes, writer)
            filesCutSave.append(outputPdf)
    except Exception as e:
        st.error(f"Ocorreu um erro: {e} - página {i+1}", icon='🛑')
    return filesCutSave    

def createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, rotate):
    numPgOne -= 1    
    inputPdf = docPdf
    name, ext = os.path.splitext(namePdf)
    typeSeq = st.session_state[listKeys[5]]
    if st.session_state[listKeys[5]] == 0:
        listSel = [pg for pg in range(numPgOne, numPgTwo)]
        outputPdf = f'{name}_{numPgOne + 1}_{numPgTwo}.pdf'    
    else:
        listSel = seqPages(numPgOne, numPgTwo)
        outputPdf = f'{name}_{numPgOne + 1}_{numPgTwo}_.pdf'
    docPdf.select(listSel)
    docPdf.save(outputPdf)
    if rotate:
        outputPdf = rotatePdf(outputPdf, index) 
    docPdf.close()
    return outputPdf   

def addWatermark(inputPdf, valMark):
    doc = pymupdf.open(inputPdf)
    for page_num in range(doc.page_count):
        page = doc[page_num]
        page_rect = page.rect
        x = page_rect.width/6.5
        y = page_rect.height * 0.95
        page.insert_text(
            (x, y),  
            valMark,  
            fontsize=25,
            color=(0.7, 0.7, 0.4),  
            rotate=0,  
            fill_opacity=0.3,
            stroke_opacity=0.3
        )
    doc.save(inputPdf, incremental=True, encryption=0)
    doc.close()
    return inputPdf

def selPdfMark(docPdf, numPgOne, numPgTwo, namePdf, index, valMark):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, True)
    pdfMark = addWatermark(outputPdf, valMark)
    downPdfUnique(pdfMark, numPgOne, numPgTwo, namePdf)           
    
def selPgsSize(docPdf, numPgOne, numPgTwo, namePdf, index, sizeMax):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, True)
    inputPdf = outputPdf
    sizeMaxStr = str(sizeMax).replace('.', '_')
    sizeSplit = sizeMaxStr.split('_')
    try:
        numOne = sizeSplit[0]
        numTwo = sizeSplit[1][:2]
        if numTwo.strip() == '00':
             numTwo = ''
        sizeMaxStr = numOne + '_' + numTwo 
    except:
        pass
    outputBase = f'{os.path.splitext(inputPdf)[0]}_divisão_{sizeMaxStr}_Mb__parte_'
    filesCutSave = divideBySize(inputPdf, sizeMax, outputBase)
    downloadExt(filesCutSave, namePdf, numPgOne, numPgTwo, 'pedaços')

@st.cache_data
def extractTables(filePdf):
    docPdf = pymupdf.open(filePdf)
    AllTable = []
    for page in docPdf:
        tabs = page.find_tables()
        for t, tab in enumerate(tabs):
            listaTable = tab.extract()
            for lista in listaTable:
                AllTable.append(lista)
    return AllTable
    
def selImgUrlsPgs(docPdf, numPgOne, numPgTwo, namePdf, mode, index):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, False)
    filesImg = extractImgs(outputPdf)
    downloadExt(filesImg, namePdf, numPgOne, numPgTwo, 'imagens')
    
@st.cache_data   
def imagesConvert(filePdf):
    docPdf = pymupdf.open(filePdf)
    nPages = len(docPdf)
    listImgs = []
    imgs = st.session_state[keyImgs] 
    for img in imgs:
        for pg in range(nPages):
            page = docPdf.load_page(pg)
            pix = page.get_pixmap()
            fileImg = f'imagem_{pg + 1}{img}'
            pix.save(fileImg)
            listImgs.append(fileImg)
    return listImgs  

@st.cache_data 
def docxConvert(filePdf):
    name = os.path.splitext(filePdf)[0]
    docs = st.session_state[keyDocs]
    listDocs = []
    for doc in docs:
        fileDoc = f'{name}{doc}'
        if doc.lower() == '.html': 
            text = extractText(filePdf, 2)
            with open(fileDoc, "w") as file:
                file.write(text)
            listDocs.append(fileDoc)
        else:
            try:
                cv = Converter(filePdf)
                cv.convert(fileDoc, start=0, end=None)
                cv.close()
                listDocs.append(fileDoc)
            except: 
                pass
    return listDocs

@st.cache_data 
def tableConvert(filePdf):   
    name = os.path.splitext(filePdf)[0]
    tables = st.session_state[keyDocs]
    listTables = []
    for table in tables:
        allTables = extractTables(filePdf)
        if len(allTables) > 0:
            fileTable = f'{name}{table}'
            workbook = xlsxwriter.Workbook(fileTable)
            worksheet = workbook.add_worksheet('aba_dados')
            for rowNum, rowData in enumerate(allTables):
                worksheet.write_row(rowNum, 0, rowData)
            workbook.close()
            listTables.append(fileTable)    
    return listTables
     
@st.cache_data 
def ppTxConvert(filePdf):
    docPdf = pymupdf.open(filePdf)
    baseName = os.path.basename(filePdf)
    name, ext = os.path.splitext(baseName)
    slides = st.session_state[keySlides]
    listSlides = []
    for slide in slides:
        fileSlide = f'{name}{slide}'
        dictAllTexts = {}
        for pg, page in enumerate(docPdf):
            nPg = pg + 1
            text = page.get_text()
            dictAllTexts.setdefault(nPg, '')
            dictAllTexts[nPg] += f'{text}\n'
        wrapper = textwrap.TextWrapper(width=75)
        p = Presentation()
        contPg = 0
        for dctAll, texts in dictAllTexts.items():
            textSplit = [txt.strip() for txt in texts.split('\n') if len(txt.strip()) != 0]
            textAdd = ''
            contSeg = 0 
            for tx, text in enumerate(textSplit):
                if tx%14 == 0 and tx != 0:
                    contSeg += 1
                    s = p.slides.add_slide(p.slide_layouts[5])
                    titlePara = s.shapes.title.text_frame.paragraphs[0]
                    titlePara.font.name = "Times New Roman"
                    titlePara.font.size = Pt(18)
                    titlePara.text = f'Arquivo {name} - página {dctAll} - segmento {contSeg}'
                    txt_box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1))
                    txt_frame = txt_box.text_frame
                    n = txt_frame.add_paragraph()
                    string = wrapper.fill(text=textAdd)
                    n.text = string
                    n.alignment = PP_ALIGN.JUSTIFY
                    textAdd = ''
                    textAdd += text + '\n'
                else:
                    textAdd += text + '\n'
            s = p.slides.add_slide(p.slide_layouts[5])
            titlePara = s.shapes.title.text_frame.paragraphs[0]
            titlePara.font.name = "Times New Roman"
            titlePara.font.size = Pt(18)
            titlePara.text = f'Arquivo {name} - página {dctAll} - segmento {contSeg}'
            txt_box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1))
            txt_frame = txt_box.text_frame
            n = txt_frame.add_paragraph()
            string = wrapper.fill(text=textAdd)
            n.text = string
            n.alignment = PP_ALIGN.JUSTIFY
        p.save(fileSlide)
        listSlides.append(fileSlide)
    return listSlides

@st.cache_data   
def createImgQrCode():
    fileImg = 'myContact.png'
    valuesQrcode = []
    for k, key in enumerate(qrCodeKeys):
        valueState = st.session_state[key]
        if type(valueState) == tuple:
            valueState = valueState[0]
        valuesQrcode.append(valueState.strip())
    qrcode = helpers.make_mecard(name=unidecode.unidecode(valuesQrcode[0]), 
                                 phone=valuesQrcode[1], 
                                 email=valuesQrcode[2])
    qrcode.save(fileImg, scale=1)
    return fileImg    
 
@st.cache_data  
def removeAllImg(filePdf):
    name, ext = os.path.splitext(filePdf)
    outputPdf = name + f'_without_img{ext}'
    docPdf = pymupdf.open(filePdf)
    for page in docPdf:       
        imgList = page.get_images()
        try:
            for img in imgList:
                page.delete_image(img[0])
        except:
            pass
    docPdf.save(outputPdf)
    return outputPdf
    
@st.cache_data  
def removeAllWords(filePdf):
    name, ext = os.path.splitext(filePdf)
    outputPdf = name + f'_without_words{ext}'
    docPdf = pymupdf.open(filePdf)
    textSearch = st.session_state[keyWords][0]
    for pageNum in range(docPdf.page_count):
        page = docPdf.load_page(pageNum)
        textInstances = page.search_for(textSearch)
        for inst in textInstances:
            rect = pymupdf.Rect(inst)
            widget = page.add_redact_annot(rect)
            page.apply_redactions()
        docPdf.save(outputPdf) 
    docPdf.close()
    return outputPdf

@st.cache_data   
def removeAllMark(filePdf):
    name, ext = os.path.splitext(filePdf)
    outputPdf = name + f'_without_mark{ext}'
    process_pdf(filePdf, outputPdf)
    return outputPdf

@st.cache_data   
def lockAllPages(filePdf):
    docPdf = pymupdf.open(filePdf)
    name, ext = os.path.splitext(filePdf)
    outputPdf = name + f'_lock{ext}'
    docPdf.save(outputPdf, 
                encryption=pymupdf.PDF_ENCRYPT_AES_256, 
                user_pw=st.session_state[keyWords][1])
    return outputPdf    

def unLockAllPages(docPdf, namePdf):
    name, ext = os.path.splitext(namePdf)
    outputPdf = name + f'_unlock{ext}'
    rc = docPdf.authenticate(st.session_state[keyWords][1])
    docPdf.save(outputPdf, 
                encryption=pymupdf.PDF_ENCRYPT_NONE)
    return outputPdf
    
def selPdfToAll(docPdf, numPgOne, numPgTwo, namePdf, index, rotate, sufix):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, rotate)
    if sufix.find('img') >= 0:
        listFiles = imagesConvert(outputPdf)
        obj = 'imagem'
    elif sufix.find('doc') >= 0:
        listFiles = docxConvert(outputPdf)
        obj = 'documento'
    elif sufix.find('table') >= 0:
        listFiles = tableConvert(outputPdf)
        obj = 'tabela'
    elif sufix.find('slide') >= 0:
        listFiles = ppTxConvert(outputPdf)
        obj = 'slide'
    if len(listFiles) > 0:
        downloadExt(listFiles, namePdf, numPgOne, numPgTwo, sufix)
    else:
        strEmpty = f'😢 Conversão em "{obj}" fracassada para o arquivo "{namePdf}", intervalo de páginas {numPgOne}-{numPgTwo}!'
        strEmpty += 'Verifique se há necessidade de OCR (reconhecimento óptico de caracteres).'
        config(strEmpty)   
    
def selPdfToQrcode(docPdf, numPgOne, numPgTwo, namePdf, index):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, True)
    fileImg = createImgQrCode()
    filePdf = insertImgPdf(outputPdf, fileImg)
    with open(filePdf, "rb") as file:
        PDFbyte = file.read()
    mensResult(0, 1, 'pdf', PDFbyte, filePdf)   

def selPdfRemoveImg(docPdf, numPgOne, numPgTwo, namePdf, index):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, True)
    filePdf = removeAllImg(outputPdf)
    with open(filePdf, "rb") as file:
        PDFbyte = file.read()
    mensResult(0, 1, 'pdf', PDFbyte, filePdf)   
    
def selPdfRemoveWords(docPdf, numPgOne, numPgTwo, namePdf, index):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, True)
    filePdf = removeAllWords(outputPdf)
    with open(filePdf, "rb") as file:
        PDFbyte = file.read()
    mensResult(0, 1, 'pdf', PDFbyte, filePdf)  
    
def selPdfRemoveMark(docPdf, numPgOne, numPgTwo, namePdf, index):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, True)
    filePdf = removeAllMark(outputPdf)
    with open(filePdf, "rb") as file:
        PDFbyte = file.read()
    mensResult(0, 1, 'pdf', PDFbyte, filePdf)
    
def selPdfLockPdf(docPdf, numPgOne, numPgTwo, namePdf, index):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, True)    
    filePdf = lockAllPages(outputPdf)
    with open(filePdf, "rb") as file:
        PDFbyte = file.read()
    mensResult(0, 1, 'pdf', PDFbyte, filePdf)
    
def selPdfUnLockPdf(docPdf, numPgOne, numPgTwo, namePdf, index):
    filePdf = unLockAllPages(docPdf, namePdf)
    with open(filePdf, "rb") as file:
        PDFbyte = file.read()
    mensResult(0, 1, 'pdf', PDFbyte, filePdf)

def insertImgPdf(filePdf, imgFile):
    baseName = os.path.basename(filePdf)
    name, ext = os.path.splitext(baseName)
    newPdf = f'{name}_img.pdf'
    headY = 720
    docPdf = pymupdf.open(filePdf)
    for pg in range(len(docPdf)):
        page = docPdf.load_page(pg)
        rect = page.rect  
        img = pymupdf.open(imgFile)
        img_rect = img.load_page(0).rect
        x0 = (rect.width - img_rect.width) / 13  
        y0 = rect.height - headY / 11            
        page.insert_image((x0, y0, x0 + img_rect.width, y0 + img_rect.height), filename=imgFile)             
    docPdf.save(newPdf)
    return newPdf    
    
def selTxtUrlPgs(docPdf, numPgOne, numPgTwo, namePdf, mode, index):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, False)
    if mode == 0:
        text = extractText(outputPdf, mode)
        strLabel = "Download_text"
        outputTxt = f'{namePdf}_{numPgOne}_{numPgTwo}_text.txt'
        strEmpty = f'😢 Extração de "texto" fracassada para o arquivo "{namePdf}", intervalo de páginas {numPgOne}-{numPgTwo}!'
        strEmpty += '\nVerifique se há necessidade de OCR (reconhecimento óptico de caracteres).'  
    elif mode == 1:
        text = extractUrls(outputPdf)
        strLabel = "Download_urls"
        outputTxt = f'{namePdf}_{numPgOne}_{numPgTwo}_urls.txt'
        strEmpty = f'😢 Extração de "URL" fracassada para o arquivo "{namePdf}", intervalo de páginas {numPgOne}-{numPgTwo}!' 
        strEmpty += '\nVerifique se há necessidade de OCR (reconhecimento óptico de caracteres).'
    else: 
        text = extractText(outputPdf, mode)
        strLabel = "Download_html"
        outputTxt = f'{namePdf}_{numPgOne}_{numPgTwo}_hmtl.txt'
        strEmpty = f'😢 Extração de texto em formato "HTML" fracassada para o arquivo "{namePdf}", intervalo de páginas {numPgOne}-{numPgTwo}!' 
        strEmpty += '\nVerifique se há necessidade de OCR (reconhecimento óptico de caracteres).'
    if len(text.strip()) > 0:
        mensResult(2, 1, 'txt', text, outputTxt)        
    else:
        config(strEmpty)    
                                          
def selDelPgs(docPdf, numPgOne, numPgTwo, namePdf, mode, index):
    numPgOne -= 1
    inputPdf = docPdf
    name, ext = os.path.splitext(namePdf)
    listPgs = seqPages(numPgOne, numPgTwo)
    if mode == 0:
        outputPdf = f'{name}_sel_{numPgOne + 1}_{numPgTwo}{ext}'
        listSel = [pg for pg in range(numPgOne, numPgTwo) if pg in listPgs]
    else:
        numPages = inputPdf.page_count
        outputPdf = f'{name}_del_{numPgOne + 1}_{numPgTwo}{ext}'
        listSel = [pg for pg in range(numPages) if pg not in range(numPgOne, numPgTwo)]
        listPlus = [pg for pg in range(numPgOne, numPgTwo) if pg not in listPgs]
        listSel = listPlus + listSel
    docPdf.select(listSel)
    docPdf.save(outputPdf)
    docPdf.close()
    if index != 4:
        outputPdf = rotatePdf(outputPdf, index)       
    downPdfUnique(outputPdf, numPgOne, numPgTwo, namePdf) 
                       
def downPdfUnique(outputPdf, numPgOne, numPgTwo, namePdf):
    with open(outputPdf, "rb") as pdf_file:
        PDFbyte = pdf_file.read()
    if len(PDFbyte) > 0:
        mensResult(0, 1, 'pdf', PDFbyte, outputPdf)
        
def extractPgs(docPdf, numPgOne, numPgTwo, mode, namePdf, index):
    numPgOne -= 1
    filesPdf = [docPdf]
    filesRead = []  
    name, ext = os.path.splitext(namePdf)
    for file in filesPdf:
        diffPg = abs(numPgTwo - numPgOne)
        minPg = min([numPgOne, numPgTwo])
        listPg = [pg for pg in range(minPg, diffPg)]  
        if st.session_state[listKeys[5]] != -1:
            for p, pageNum in enumerate(listPg):
                inputPdf = file
                outputPdf = f'{name}_{pageNum + 1}.pdf'
                newPdf = pymupdf.open()
                newPdf.insert_pdf(inputPdf, from_page=pageNum, to_page=pageNum)
                newPdf.save(outputPdf)
                if index != 4:
                    outputPdf = rotatePdf(outputPdf, index) 
                filesRead.append(outputPdf)
                newPdf.close()
        else:
            listPgBlock = [pg for pg in range(minPg, diffPg)] 
            sizeSegment = st.session_state[listKeys[6]]
            listAllBlocks = list(itertools.batched(listPg, sizeSegment))
            for listPg in listAllBlocks:
                inputPdf = file
                pageIni = listPg[0] + 1
                pageFinal = listPg[-1] + 1
                outputPdf = f'{name}_pgs_{pageIni}_{pageFinal}.pdf'
                newPdf = pymupdf.open()
                newPdf.insert_pdf(inputPdf, from_page=pageIni, to_page=pageFinal)
                newPdf.save(outputPdf)
                if index != 4:
                    outputPdf = rotatePdf(outputPdf, index) 
                filesRead.append(outputPdf)
                newPdf.close()                
    downloadExt(filesRead, namePdf, numPgOne, numPgTwo, 'páginas')

def exibeInfo(docPdf):
    infoDictKeys = {'Metadado': [], 'Informação': []}
    @st.dialog(' ')
    def config():
        trace = '👎'
        nPgs = docPdf.page_count
        size = round(uploadPdf.size/1024, 2)
        if size > 1024:
            size /= 1024
            size = round(size, 2)
            unit = 'MB'
        else:
            unit = 'KB'
        typFile = uploadPdf.type
        dirty = docPdf.is_dirty
        pdfYes = docPdf.is_pdf
        close = docPdf.is_closed
        formPdf = docPdf.is_form_pdf
        encry = docPdf.is_encrypted
        pdfMeta = docPdf.metadata
        infoDictKeys['Metadado'].append('🗄️ tamanho')
        infoDictKeys['Informação'].append(f'{size}{unit}')
        infoDictKeys['Metadado'].append('📄️ páginas')
        infoDictKeys['Informação'].append(nPgs)
        dictKeys = {'creator': '💂 criador', 'producer': '🔴 responsável', 'creationDate': '📅 dia de criação', 
                    'modDate': '🕰️ dia de modificação', 'title': '#️⃣  título', 'author': '📕 autor', 'format': '⏹️ formato',
                    'subject': '🖊️ assunto', 'keywords': '#️⃣  palavras-chave', 'encryption': '🔑 criptografia'}
        listDictKeys = list(dictKeys.keys())
        keys = [key for key in listDictKeys]
        for k, key in enumerate(keys):
            valueKey = dictKeys[key]
            metaKey = pdfMeta[key]
            if metaKey is None:
                metaKey = trace
            else:
                if len(metaKey.strip()) == 0:
                    metaKey = trace
            if k in [2, 3]:
                metaKey = configDate(metaKey) 
            #infoDictKeys = {'Metadado': [], 'Informação': []}
            infoDictKeys['Metadado'].append(dictKeys[key])
            infoDictKeys['Informação'].append(metaKey)
        df = pd.DataFrame(infoDictKeys)
        st.dataframe(data=df, use_container_width=True, hide_index=True, 
                     column_config={'Metadado': st.column_config.TextColumn(width='medium', 
                                                                            help=f'Metadados do arquivo'), 
                                    'Informação': st.column_config.TextColumn(width='large', 
                                                                            help=f'Valor associado ao metadado')})
    config()
    
def configDate(datePdf):
    try:
        dateSplit = datePdf.split(':')
        dateStr = dateSplit[1][:14]
        year = dateStr[:4]
        month = dateStr[4:6]
        day = dateStr[6:8]
        hour = dateStr[8:10]
        minute = dateStr[10:12]
        second = dateStr[12:]
        dateStr = f'{day}/{month}/{year}, {hour}h{minute}min{second}s'
        return dateStr
    except:
        return datePdf

@st.cache_data 
def validateEmail(email):
    return is_valid_email(email)

@st.cache_data  
def validatePhone(phone):
    newPhone = remove_symbols_phone(phone)
    return is_valid_phone(newPhone)
    
def exibeQrCode():
    @st.dialog('Dados')
    def config():
        nameUser = st.text_input(label='Nome', help='Informe nome e sobrenome. A acentuação será desconsiderada.',
                                 key=qrCodeKeys[0], placeholder=valuesReserve[0], value='')
        phoneUser = st.text_input(label='Telefone', help='Informe código de área (DDD) e demais dígitos. O uso de separadores é opcional.',
                                  key=qrCodeKeys[1], placeholder=valuesReserve[1], value=''), 
        emailUser = st.text_input(label='E-mail', help="Informe o correio eletrônico completo. O uso de '@' e ponto ('.') é obrigatório.",
                                  key=qrCodeKeys[2], placeholder=valuesReserve[2], value='')
        if len(phoneUser[0].strip()) > 0:
            if not validatePhone(phoneUser[0]):
                st.error(f"O telefone '{phoneUser[0]}' não é válido! Tente de novo!")
                time.sleep(2)
        if emailUser:
            if not validateEmail(emailUser):
                st.error(f"O e-mail '{emailUser}' não é válido! Tente de novo!")
                time.sleep(2)
        buttReturn = st.button('retornar')
        if buttReturn:
            for key in qrCodeKeys:
                del st.session_state[key]
            st.session_state[qrCodeKeys[0]] = nameUser
            st.session_state[qrCodeKeys[1]] = phoneUser[0]
            st.session_state[qrCodeKeys[2]] = emailUser
            st.rerun()
    config()

def exibeWord():
    @st.dialog(' ') 
    def config():
        words = ['Exclusão', 'Bloqueio/desbloqueio', "Marca d'água"]
        optWord = st.radio(label='Opções de Texto', options=words, index=None, horizontal=True)
        if optWord:
            ind = words.index(optWord)
            if ind == 1:
                typeInput = 'password'
                wordInput = 'Digite a senha para bloquear/desbloquear o arquivo.'
            else:
                typeInput = 'default'
                wordInput = 'Digite a palavra a ser deletada. Maiúsculas/minúsculas são ignoradas, mas a acentuação será levada em conta.'
            optWord = st.text_input(label=f'Digite a {optWord.lower()}.', key=keyWords[ind], placeholder='', 
                                    value='', type=typeInput, help=wordInput)
            buttReturn = st.button('retornar')        
            if buttReturn:
                del st.session_state[keyWords]
                st.session_state[keyWords] = ['', '', '']
                st.session_state[keyWords][ind] = optWord
                st.rerun()            
    config()
                
@st.dialog(' ')
def config(str):
    st.markdown(str)  
    
@st.dialog(' ')
def configTwo(str):
    st.markdown(str)  

@st.dialog(' ')
def configSucess(str, icon):
    st.markdown(str)   
    colMens, colDown = st.columns([8, 2])    
    
@st.dialog(' ')
def windowAdd(numPgOne, numPgTwo):
    opts = ['Seleção de páginas', 'Bloco de páginas']
    segm = st.pills(label='Formatos de saída', options=opts, selection_mode='single', 
                                 default=None, width='stretch')
    try:
        if segm.strip() == opts[1]: 
            valueNum = st.number_input(label=f'Número de páginas por bloco', min_value=numPgOne, max_value=numPgTwo)
            del st.session_state[listKeys[6]]
            st.session_state[listKeys[6]] = valueNum
            del st.session_state[listKeys[5]]
            st.session_state[listKeys[5]] = -1
        elif segm.strip() == opts[0]:
            colMode, colValue = st.columns([5, 3])
            selModel = colMode.selectbox(label=f'Páginas no intervalo de :blue[**{numPgOne}**] a :blue[**{numPgTwo}**]', 
                                         options=optionsSel)
            if selModel:
                if selModel == optionsSel[-1]:
                    valueNum = colValue.number_input(label=f'Múltiplo', min_value=numPgOne, max_value=numPgTwo) 
                    del st.session_state[listKeys[6]]
                    st.session_state[listKeys[6]] = valueNum            
                del st.session_state[listKeys[5]]
                st.session_state[listKeys[5]] = optionsSel.index(selModel)
    except:
        pass
    if st.button('retornar'):
        st.rerun()
        
@st.dialog(' ')
def windowDocsImgs(keys, mode):
    match mode:
        case 0:
            docFormats = ['.csv', '.ods', '.xls', '.xlsx']
            listSpace = [12, 9]
        case 1:
            docFormats = ['.doc', '.docx', '.odt', '.rtf', '.html']
            listSpace = [10, 6]
        case 2: 
            docFormats = ['.jpg', '.jpeg', '.png', '.pnm']
            listSpace = [12, 9]
        case _:
            docFormats = ['.odp', '.ppt', '.pptx']
            listSpace = [16, 12]
    colSeg, colMark = st.columns(listSpace, vertical_alignment='bottom', width='stretch')
    if mode in [0, 2]: 
        segms = colSeg.segmented_control(label='Formatos de saída', options=sorted(docFormats), selection_mode='multi', 
                                         default=None)
    else:
        segms = colSeg.pills(label='Formatos de saída', options=docFormats, selection_mode='multi', 
                             default=None)
    nSegms = len(segms)
    if nSegms > 0:
        colMark.markdown('')
        if nSegms == 1:
            expr = 'formato'
        else:
            expr = 'formatos'    
        colMark.markdown(f'{nSegms} {expr}:\n{segms}')
    else:
        textIni = 'Nada selecionado!'
        colMark.code(textIni.ljust(len(textIni), ' '), language='Python')
    if st.button('retornar'):
        del st.session_state[keys]
        st.session_state[keys] = []
        st.session_state[keys] = segms
        st.rerun()
       
def iniFinally(mode):
    if mode == 0:
        for key in listKeys:
            if key not in st.session_state:
                try:
                    st.session_state[key] = dictKeys[key]
                except:
                    pass        
    else:
        try:
            for key in listKeys:
                del st.session_state[key]
        except:
            pass  
        iniFinally(0)
        st.rerun()
        
def seqPages(numPgOne, numPgTwo):
    valNum = st.session_state[listKeys[5]] 
    listPgs = [pg for pg in range(numPgOne, numPgTwo)]
    match valNum:
        case 1:
            listPgs = [pg for pg in range(numPgOne, numPgTwo) if (pg+1)%2==0]
        case 2:
            listPgs = [pg for pg in range(numPgOne, numPgTwo) if (pg+1)%2==1]
        case 3:
            listPgs = [pg for pg in range(numPgOne, numPgTwo)]
        case 4:
            mult = st.session_state[listKeys[6]]
            listPgs = [pg for pg in range(numPgOne, numPgTwo) if (pg+1)%mult==0]        
    return listPgs  

def upDownScroll(w):
    label = icons[w][0]
    jsCode = f"""
        <script>
            var lastNameInput = window.parent.document.querySelector('input[aria-label={label}]');
            lastNameInput.focus();
        </script>
    """
    components.html(jsCode)
    
def main():
    global uploadPdf
    global valMx
    global sufix
    sufix = ['']
    with st.container(border=None, key='contOne', vertical_alignment='top'):
        uploadPdf = st.file_uploader('Selecionar arquivo PDF', 
                                     type=['pdf'], 
                                     accept_multiple_files=False,
                                     label_visibility="collapsed")
        if uploadPdf is not None:
            number = 4
            pdfName = uploadPdf.name
            docPdf = pymupdf.open(stream=uploadPdf.read(), filetype="pdf")
            valMx = docPdf.page_count 
            valMxSize = round(uploadPdf.size/(1024**2), 2)
            sizeColsDate = [1.2, 1.2, 1.4, 1.8]
            lenColsDate = len(sizeColsDate)
            sizeColsSuppOne = [1 for n in range(number)]
            sizeColsSuppTwo = sizeColsSuppOne
            lenColsSupp = len(sizeColsSuppOne) + len(sizeColsSuppTwo)
            sizeColsText = sizeColsSuppOne
            lenColsText = len(sizeColsText)
            sizeColsMultOne = [1 for n in range(number)]
            sizeColsMultTwo = [1 for n in range(number)]
            sizeColsMultThree = [1 for n in range(number)]
            lenColsMult = len(sizeColsMultOne) + len(sizeColsMultTwo) + len(sizeColsMultThree)
            sizeColsFormat = [1 for n in range(number)]
            lenColsFormat = len(sizeColsFormat)
            with st.container(border=4, key='contTwo'):
                colEmptyOne, colDate, colEmptyTwo = st.columns(spec=3, vertical_alignment='bottom', 
                                                               width='stretch')
                colDate.markdown(f':material/cards: :blue[**Parâmetros básicos**] (:red[**{lenColsDate}**])', unsafe_allow_html=True, 
                                 help='Exibe as opções básicas de paginação, ângulo de rotação e de divisão por tamanho.')
                colPgOne, colPgTwo, colSize, colSlider = st.columns(sizeColsDate, vertical_alignment='bottom', 
                                                                    width='stretch')
                numPgOne = colPgOne.number_input(label='Página inicial  (:red[**1**])', key=listKeys[0], 
                                                 min_value=1, max_value=valMx, 
                                                 help=f'Digite, incremente ou decremente um número entre "1" e "{valMx}".')
                numPgTwo = colPgTwo.number_input(label=f'Página final  (:red[**{valMx}**])', key=listKeys[1], 
                                                 min_value=1, max_value=valMx, 
                                                 help=f'Digite, incremente ou decremente um número até o máximo de "{valMx}".')
                valPgSize = colSize.number_input(label='Tamanho para divisão (:red[**MB**])', key=listKeys[3], 
                                                 min_value=dictKeys[listKeys[3]], step=dictKeys[listKeys[3]],  
                                                 max_value=valMxSize, 
                                                 help='Digite, incremente ou decremente o tamanho de cada fração do arquivo.')
                valPgAngle = colSlider.select_slider(label='Ângulo de rotação', options=valAngles, 
                                                     key=listKeys[2], 
                                                     help='Escolha o ângulo de rotação deslizando o botão para a esquerda ou direita.')
            with st.container(border=6, key='webDown'):                
                colIcoFinal, colIcoBasic, colIcoExtract, colIcoOther, colIcoConv, colButtFinal, colButtClear = st.columns([3, 10, 10, 10, 10, 10, 10], 
                                                                                                              vertical_alignment='center', 
                                                                                                              width='stretch')
                textIcoFinal = colIcoFinal.text_input(icons[0][0], icon=icons[0][1], width=40)
                buttBottomBasic = colIcoBasic.button(label=dictButts[keysButts[31]][0], key=keysButts[31], 
                                                     use_container_width=True, icon=dictButts[keysButts[31]][1], 
                                                     help=dictButts[keysButts[31]][-1]) 
                buttBottomExtract = colIcoExtract.button(label=dictButts[keysButts[32]][0], key=keysButts[32], 
                                                     use_container_width=True, icon=dictButts[keysButts[32]][1], 
                                                     help=dictButts[keysButts[32]][-1])
                buttBottomOther = colIcoOther.button(label=dictButts[keysButts[33]][0], key=keysButts[33], 
                                                     use_container_width=True, icon=dictButts[keysButts[33]][1], 
                                                     help=dictButts[keysButts[33]][-1])                
                buttBottomConv = colIcoConv.button(label=dictButts[keysButts[34]][0], key=keysButts[34], 
                                                     use_container_width=True, icon=dictButts[keysButts[34]][1], 
                                                     help=dictButts[keysButts[34]][-1])                
                buttBottomWeb = colButtFinal.button(label=dictButts[keysButts[29]][0], key=keysButts[29], 
                                                    use_container_width=True, icon=dictButts[keysButts[29]][1], 
                                                    help=dictButts[keysButts[29]][-1]) 
                buttPgClear = colButtClear.button(label=dictButts[keysButts[4]][0], key=keysButts[4], 
                                                  use_container_width=True, icon=dictButts[keysButts[4]][1], 
                                                  help=dictButts[keysButts[4]][-1]) 
            with st.container(border=4, key='contThree'):
                colBasic, colSupport, colEmptyOnly = st.columns(spec=3, vertical_alignment='center', 
                                                                  width='stretch')
                textIcoBasic = colBasic.text_input(icons[1][0], icon=icons[1][1], width=40)
                colSupport.markdown(f':material/cards: :blue[**Telas de apoio**] (:red[**{lenColsSupp}**])', unsafe_allow_html=True, 
                                    help='Exibe as opções para as telas de apoio.')                                                        
                colPgs, colWords, colptRotate, colOptPlans = st.columns(sizeColsSuppOne, vertical_alignment='bottom', 
                                                                        width='stretch')                
                colOptDocs, colOptImgs, colOptSlides, colPerson = st.columns(sizeColsSuppTwo, vertical_alignment='bottom', 
                                                                             width='stretch')
                buttToPages = colPgs.button(label=dictButts[keysButts[15]][0], use_container_width=True, 
                                            icon=dictButts[keysButts[15]][1], key=keysButts[15], 
                                            help=dictButts[keysButts[15]][-1])
                buttOptWords = colWords.button(label=dictButts[keysButts[19]][0], use_container_width=True, 
                                               icon=dictButts[keysButts[19]][1], key=keysButts[19], 
                                               help=dictButts[keysButts[19]][-1]) 
                buttOptRotate = colptRotate.button(label=dictButts[keysButts[27]][0], use_container_width=True, 
                                                   icon=dictButts[keysButts[27]][1], key=keysButts[27], 
                                                   help=dictButts[keysButts[27]][-1])
                buttOptPlans = colOptPlans.button(label=dictButts[keysButts[23]][0], use_container_width=True, 
                                                  icon=dictButts[keysButts[23]][1], key=keysButts[23], 
                                                  help=dictButts[keysButts[23]][-1]) 
                buttOptDocs = colOptDocs.button(label=dictButts[keysButts[25]][0], use_container_width=True, 
                                                icon=dictButts[keysButts[25]][1], key=keysButts[25], 
                                                help=dictButts[keysButts[25]][-1])
                buttOptImgs = colOptImgs.button(label=dictButts[keysButts[24]][0], use_container_width=True, 
                                                icon=dictButts[keysButts[24]][1], key=keysButts[24], 
                                                help=dictButts[keysButts[24]][-1]) 
                buttOptSlides = colOptSlides.button(label=dictButts[keysButts[26]][0], use_container_width=True, 
                                                    icon=dictButts[keysButts[26]][1], key=keysButts[26], 
                                                    help=dictButts[keysButts[26]][-1]) 
                buttPerson = colPerson.button(label=dictButts[keysButts[16]][0], use_container_width=True, 
                                                  icon=dictButts[keysButts[16]][1], key=keysButts[16], 
                                                  help=dictButts[keysButts[16]][-1])
            with st.container(border=4, key='contFour'):
                colBasic, colExtract, colEmptyOnly = st.columns(spec=3, vertical_alignment='center', 
                                                                  width='stretch')
                textIcoExtract = colBasic.text_input(icons[2][0], icon=icons[2][1], width=40)
                colExtract.markdown(f':material/cards: :blue[**Extração de objetos**] (:red[**{lenColsText}**])', unsafe_allow_html=True, 
                                    help='Exibe as opções para extração de texto, imagem, URL e código HTML.') 
                colButtTxt, colButtImg, colButtUrl, colButtHtml= st.columns(sizeColsText, vertical_alignment='bottom', 
                                                                 width='stretch')
                buttPgTxt = colButtTxt.button(label=dictButts[keysButts[1]][0], key=keysButts[1], 
                                              use_container_width=True, icon=dictButts[keysButts[1]][1], 
                                              help=dictButts[keysButts[1]][-1])
                buttPdfImg = colButtImg.button(label=dictButts[keysButts[6]][0], key=keysButts[6], 
                                               use_container_width=True, icon=dictButts[keysButts[6]][1], 
                                               help=dictButts[keysButts[6]][-1])
                buttPdfUrl = colButtUrl.button(label=dictButts[keysButts[5]][0], key=keysButts[5], 
                                               use_container_width=True, icon=dictButts[keysButts[5]][1], 
                                               help=dictButts[keysButts[5]][-1])
                buttPdfHtml = colButtHtml.button(label=dictButts[keysButts[28]][0], key=keysButts[28], 
                                                   use_container_width=True, icon=dictButts[keysButts[28]][1], 
                                                   help=dictButts[keysButts[28]][-1])
            with st.container(border=4, key='contFive'):
                colBasic, colMult, colEmptyOnly = st.columns(spec=3, vertical_alignment='center', 
                                                                  width='stretch')
                textIcoOther = colBasic.text_input(icons[3][0], icon=icons[3][1], width=40)
                colMult.markdown(f':material/cards: :blue[**Outras operações**] (:red[**{lenColsMult}**])', unsafe_allow_html=True, 
                                  help='Exibe diversas opções de tratamento/manipulação de PDF.') 
                colButtDiv, colButtSize, colButtSel, colCode = st.columns(sizeColsMultOne, vertical_alignment='bottom', 
                                                                          width='stretch')
                colButtMark, colRemoveMark, colCodePdf, colDecodePdf = st.columns(sizeColsMultTwo, vertical_alignment='bottom', 
                                                                                  width='stretch')
                colRemImg, colRemWrd, colButtDel, colButtInfo = st.columns(sizeColsMultThree, vertical_alignment='bottom', 
                                                                           width='stretch')
                                                                          
                buttPdfDiv = colButtDiv.button(label=dictButts[keysButts[0]][0], key=keysButts[0], 
                                              use_container_width=True, icon=dictButts[keysButts[0]][1], 
                                              help=dictButts[keysButts[0]][-1])
                buttPdfSize = colButtSize.button(label=dictButts[keysButts[7]][0], key=keysButts[7], 
                                                 use_container_width=True, icon=dictButts[keysButts[7]][1], 
                                                 help=dictButts[keysButts[7]][-1])
                buttPgSel = colButtSel.button(label=dictButts[keysButts[2]][0], key=keysButts[2], 
                                              use_container_width=True, icon=dictButts[keysButts[2]][1], 
                                              help=dictButts[keysButts[2]][-1])
                buttQrcode =  colCode.button(label=dictButts[keysButts[14]][0], key=keysButts[14], 
                                             use_container_width=True, icon=dictButts[keysButts[14]][1], 
                                             help=dictButts[keysButts[14]][-1])                                             
                buttPdfMark = colButtMark.button(label=dictButts[keysButts[8]][0], key=keysButts[8], 
                                                 use_container_width=True, icon=dictButts[keysButts[8]][1], 
                                                 help=dictButts[keysButts[8]][-1])
                buttRemoveMark = colRemoveMark.button(label=dictButts[keysButts[22]][0], key=keysButts[22], 
                                                      use_container_width=True, icon=dictButts[keysButts[22]][1], 
                                                      help=dictButts[keysButts[22]][-1])
                buttCodePdf = colCodePdf.button(label=dictButts[keysButts[20]][0], key=keysButts[20], 
                                                use_container_width=True, icon=dictButts[keysButts[20]][1], 
                                                help=dictButts[keysButts[20]][-1])
                buttDecodePdf = colDecodePdf.button(label=dictButts[keysButts[21]][0], key=keysButts[21], 
                                                    use_container_width=True, icon=dictButts[keysButts[21]][1], 
                                                    help=dictButts[keysButts[21]][-1])
                buttRemoveImg = colRemImg.button(label=dictButts[keysButts[17]][0], key=keysButts[17], 
                                                 use_container_width=True, icon=dictButts[keysButts[17]][1], 
                                                 help=dictButts[keysButts[17]][-1])
                buttRemoveWords = colRemWrd.button(label=dictButts[keysButts[18]][0], key=keysButts[18], 
                                                   use_container_width=True, icon=dictButts[keysButts[18]][1], 
                                                   help=dictButts[keysButts[18]][-1])
                buttPgDel = colButtDel.button(label=dictButts[keysButts[3]][0], key=keysButts[3], 
                                              use_container_width=True, icon=dictButts[keysButts[3]][1], 
                                              help=dictButts[keysButts[3]][-1])
                buttPdfInfo =  colButtInfo.button(label=dictButts[keysButts[9]][0], key=keysButts[9], 
                                                  use_container_width=True, icon=dictButts[keysButts[9]][1], 
                                                  help=dictButts[keysButts[9]][-1])            
            with st.container(border=4, key='contSix'):
                colBasic, colMult, colEmptyOnly = st.columns(spec=3, vertical_alignment='center', 
                                                               width='stretch')
                textIcoConv = colBasic.text_input(icons[4][0], icon=icons[4][1], width=40)
                colMult.markdown(f':material/cards: :blue[**Conversão de formato**] (:red[**{lenColsFormat}**])', unsafe_allow_html=True, 
                                  help='Exibe opções de conversão de PDF em outros formatos.')
                colToTable, colToWord, colToImg, colToPower = st.columns(sizeColsFormat, vertical_alignment='bottom', 
                                                                                width='stretch')
                buttToTable = colToTable.button(label=dictButts[keysButts[10]][0], key=keysButts[10], 
                                                  use_container_width=True, icon=dictButts[keysButts[10]][1], 
                                                  help=dictButts[keysButts[10]][-1])
                buttToWord = colToWord.button(label=dictButts[keysButts[11]][0], key=keysButts[11], 
                                               use_container_width=True, icon=dictButts[keysButts[11]][1], 
                                               help=dictButts[keysButts[11]][-1])
                buttToImg = colToImg.button(label=dictButts[keysButts[12]][0], key=keysButts[12], 
                                            use_container_width=True, icon=dictButts[keysButts[12]][1], 
                                            help=dictButts[keysButts[12]][-1])
                buttToPower = colToPower.button(label=dictButts[keysButts[13]][0], key=keysButts[13], 
                                                use_container_width=True, icon=dictButts[keysButts[13]][1], 
                                                help=dictButts[keysButts[13]][-1])       
            with st.container(border=4, key='webUp'):                      
                colIcoIni, colButtIni = st.columns([1,22], vertical_alignment='center', 
                                                   width='stretch')
                textIcoIni = colIcoIni.text_input(icons[5][0], icon=icons[5][1], width=40)
                buttTopWeb = colButtIni.button(label=dictButts[keysButts[30]][0], key=keysButts[30], 
                                               use_container_width=True, icon=dictButts[keysButts[30]][1], 
                                               help=dictButts[keysButts[30]][-1])
            if numPgTwo >= numPgOne: 
                numPgIni = numPgOne
                numPgFinal = numPgTwo
            else:
                numPgIni = numPgTwo
                numPgFinal = numPgOne 
            indexAng = valAngles.index(valPgAngle)
            exprPre = f'o intervalo de páginas {numPgOne} a {numPgTwo}.' 
            if buttToPages:
                windowAdd(numPgOne, numPgTwo)
            if buttOptPlans:
               windowDocsImgs(keyTables, 0)
            if buttOptDocs: 
                windowDocsImgs(keyDocs, 1)  
            if buttOptImgs:
                windowDocsImgs(keyImgs, 2) 
            if buttOptSlides: 
                windowDocsImgs(keySlides, 3)
            if buttBottomWeb: 
               upDownScroll(5) 
            if buttTopWeb:
               upDownScroll(0) 
            if buttBottomBasic:
               upDownScroll(1)
            if buttBottomExtract:
               upDownScroll(2) 
            if buttBottomOther:
               upDownScroll(3) 
            if buttBottomConv:
               upDownScroll(4) 
            if buttPdfDiv:  
                try:
                    expr = f'{dictButts[keysButts[0]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        extractPgs(docPdf, numPgIni, numPgFinal, 0, pdfName, indexAng)
                except:
                    config(f'😢 Divisão fracassada!\n🔴 arquivo {namePdf}, intervalo de páginas {numPgOne}-{numPgTwo}!') 
            if buttPerson:
                exibeQrCode()
            if buttPgSel:
                try:
                    expr = f'{dictButts[keysButts[2]][2]} {pdfName} {exprPre}'
                    with st.spinner(expr):
                        selDelPgs(docPdf, numPgOne, numPgTwo, pdfName, 0, indexAng)
                except:
                    config(f'😢 Seleção de páginas fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttPgDel: 
                try:
                    expr = f'{dictButts[keysButts[3]][2]} {pdfName} {exprPre}'
                    with st.spinner(expr):
                        selDelPgs(docPdf, numPgOne, numPgTwo, pdfName, 1, indexAng)
                except:
                    config(f'😢 Deleção de páginas fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttPgClear: 
                del st.session_state[listKeys[5]]
                st.session_state[listKeys[5]] = 0
                iniFinally(1) 
            if buttPgTxt: 
                try:
                    expr = f'{dictButts[keysButts[1]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        selTxtUrlPgs(docPdf, numPgOne, numPgTwo, pdfName, 0, indexAng)
                except:
                    config(f'😢 Extração de texto fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttPdfImg: 
                try:     
                    expr = f'{dictButts[keysButts[6]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        sufix[0] = 'imgs'
                        selImgUrlsPgs(docPdf, numPgOne, numPgTwo, pdfName, 2, indexAng)
                except:
                    config(f'😢 Extração de imagens fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!') 
            if buttPdfUrl:
                try:
                    expr = f'{dictButts[keysButts[5]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        sufix[0] = 'urls'
                        selTxtUrlPgs(docPdf, numPgOne, numPgTwo, pdfName, 1, indexAng)
                except:
                    config(f'😢 Extração de link fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttPdfHtml:
                try:
                    expr = f'{dictButts[keysButts[28]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        sufix[0] = 'html'
                        selTxtUrlPgs(docPdf, numPgOne, numPgTwo, pdfName, 2, indexAng)
                except:
                    config(f'😢 Extração de texto em HTML fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttPdfSize:
                expr = f'{dictButts[keysButts[7]][2]} {pdfName} n{exprPre}'
                try:                
                    with st.spinner(expr):
                        selPgsSize(docPdf, numPgOne, numPgTwo, pdfName, indexAng, valPgSize)
                except:
                    configTwo(f'😢 Divisão em pedaços fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')               
            if buttPdfMark:
                try:
                    valPgMark = st.session_state[keyWords][2].strip()
                    if valPgMark.strip() == '':
                        config("😢 Nenhuma marca d'água foi digitada!\nAbra a tela e digite o texto desejado!") 
                    else:
                        expr = f'{dictButts[keysButts[8]][2]} {pdfName} n{exprPre}'
                        with st.spinner(expr):
                            selPdfMark(docPdf, numPgOne, numPgTwo, pdfName, indexAng, valPgMark)
                except:
                    config(f'😢 Marcação de páginas fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttPdfInfo:
                try:
                    expr = f'{dictButts[keysButts[9]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        exibeInfo(docPdf)
                except:
                    config(f'😢 Exibição fracassada!\n🔴 arquivo {pdfName}!')
            if buttPgClear: 
                del st.session_state[listKeys[5]]
                st.session_state[listKeys[5]] = 0
                iniFinally(1) 
            if buttToTable:
                nTables = len(st.session_state[keyTables])
                if nTables == 0:
                    config('😢 Nenhum tipo de tabela de saída foi escolhido!\nAbra a tela para realizar essa escolha!')
                else:
                    try:
                        expr = f'{dictButts[keysButts[10]][2]} {pdfName} n{exprPre}'
                        with st.spinner(expr):
                            selPdfToAll(docPdf, numPgOne, numPgTwo, pdfName, indexAng, False, 'pdf_table')          
                    except:
                        config(f'😢 Extração de tabelas fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttToWord:
                nDocs = len(st.session_state[keyDocs])
                if nDocs == 0:
                    config('😢 Nenhum tipo de documento de saída foi escolhido!\nAbra a tela para realizar essa escolha!')
                else:
                    try:
                        expr = f'{dictButts[keysButts[11]][2]} {pdfName} n{exprPre}'
                        with st.spinner(expr):
                            selPdfToAll(docPdf, numPgOne, numPgTwo, pdfName, indexAng, False, 'pdf_doc') 
                    except:
                        config(f'😢 Conversão de PDF em docx fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttToImg:
                nImgs = len(st.session_state[keyImgs])
                if nImgs == 0:
                    config('😢 Nenhum tipo de imagem foi escolhido!\nAbra a tela para realizar essa escolha!')
                else:
                    try:
                        expr = f'{dictButts[keysButts[12]][2]} {pdfName} n{exprPre}'
                        with st.spinner(expr):
                            selPdfToAll(docPdf, numPgOne, numPgTwo, pdfName, indexAng, True, 'pdf_img')
                    except: 
                        strEmpty = f'😢 Conversão de PDF em "imagem" fracassada para o arquivo "{pdfName}", intervalo de páginas {numPgOne}-{numPgTwo}!'
                        config(strEmpty)
            if buttToPower:
                nSlides = len(st.session_state[keySlides])
                if nSlides == 0:
                    config('😢 Nenhum tipo de slide foi escolhido!\nAbra a tela para realizar essa escolha!')
                else:
                    try:
                        expr = f'{dictButts[keysButts[13]][2]} {pdfName} n{exprPre}'
                        with st.spinner(expr):
                            selPdfToAll(docPdf, numPgOne, numPgTwo, pdfName, indexAng, False, 'pdf_slide')                      
                    except:
                        strEmpty = f'😢 Conversão de PDF em "slide" fracassada para o arquivo "{pdfName}", intervalo de páginas {numPgOne}-{numPgTwo}!'
                        config(strEmpty)
            if buttQrcode:
                failCode = False
                for code in qrCodeKeys: 
                    valueState = st.session_state[code]
                    if valueState.strip() == '':
                        failCode = True
                        break
                if failCode:
                    config(f'😢 Faltam dados para inserção do qrCode!\nAbra a tela e complemente-os!')
                else:
                    try:
                        expr = f'{dictButts[keysButts[14]][2]} {pdfName} n{exprPre}'
                        with st.spinner(expr):
                            selPdfToQrcode(docPdf, numPgOne, numPgTwo, pdfName, indexAng)                        
                    except:
                        config(f'😢 Inserção de QRcode fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')  
            if buttRemoveImg:
                expr = f'{dictButts[keysButts[17]][2]} {pdfName} n{exprPre}'
                try:
                    with st.spinner(expr):
                        selPdfRemoveImg(docPdf, numPgOne, numPgTwo, pdfName, indexAng)
                except:
                    config(f'😢 Remoção de imagens fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!') 
            if buttOptWords:
                exibeWord()
            if buttRemoveWords:
                try:
                    textWrite = st.session_state[keyWords][0].strip()
                    wordOk = True
                except:
                    wordOk = False
                if wordOk:
                    if len(textWrite) == 0:
                        config(f'😢 Nenhum texto foi selecionado!\nAbra a tela e digite o texto desejado!') 
                    else:
                        expr = f'{dictButts[keysButts[18]][2]} {pdfName} n{exprPre}'
                        try:
                            with st.spinner(expr):
                                selPdfRemoveWords(docPdf, numPgOne, numPgTwo, pdfName, indexAng)
                        except:
                            config(f'😢 Deleção de texto fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!') 
                else:
                    config('😢 Nenhum texto foi digitado!\nAbra a tela e digite o texto desejado!') 
            if buttCodePdf or buttDecodePdf:
                try:
                    if buttCodePdf:
                        block = ''
                    if buttDecodePdf:
                        block = 'des'
                    textWrite = st.session_state[keyWords][1].strip()
                    wordOk = True
                except:
                    wordOk = False
                if wordOk:
                    if len(textWrite) == 0:
                        config(f'😢 Nenhuma senha para {block}bloqueio foi informada!\nAbra a tela e digite o texto desejado!') 
                    else:
                        try:
                            if block == '':
                                try:
                                    expr = f'{dictButts[keysButts[20]][2]} {pdfName} n{exprPre}'
                                    with st.spinner(expr):
                                        selPdfLockPdf(docPdf, numPgOne, numPgTwo, pdfName, indexAng)
                                except:
                                    config(f'😢 {operStr} fracassado!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
                            else:
                                try:
                                    numPgTwo = valMx
                                    selPdfUnLockPdf(docPdf, numPgOne, numPgTwo, pdfName, indexAng)
                                except:
                                    config(f'😢 {operStr} fracassado!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
                        except:
                            oper = f'{block}bloqueio'
                            operStr = f'{oper.capitalize()}'
                            config(f'😢 {operStr} fracassado!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
                else:
                    config(f'😢 Nenhuma senha para {block}bloqueio foi informada!\nAbra a tela e digite o texto desejado!')                    
            if buttRemoveMark:
                expr = f'{dictButts[keysButts[22]][2]} {pdfName} n{exprPre}'
                try:
                    expr = f'{dictButts[keysButts[22]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        selPdfRemoveMark(docPdf, numPgOne, numPgTwo, pdfName, indexAng)
                except:
                    config(f"😢 Remoção de marca d'água fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!")
                        
if __name__ == '__main__':
    global dictKeys, listKeys 
    global valAngles, valComps
    global countPg, optionsSel
    global namesTeste, nameApp 
    global qrCodeKeys, valuesReserve
    global dictButts, keysButts
    global keyWords, keyDocs, keyImgs, keyTables, keySlides
    global icons
    icons = [['📃', '📌'], 
             ['🛠', '📌'],  
             ['🔧', '📌'], 
             ['🔨', '📌'], 
             ['⛏', '📌'], 
             ['🧾', '📌'], 
             ['👍', '📌']]
    nameApp = 'Ferramentas/PDF'
    valAngles = ['-360°', '-270°', '-180°', '-90°', '0°', '90°', '180°', '270°', '360°']
    optionsSel = ['', 'pares', 'não pares', 'todos', 'múltiplos de ']
    dictKeys = {'pgOne': 1, 
                'pgTwo': 1, 
                'pgAngle': valAngles[0], 
                'pgSize': 0.01, 
                'pgMark': '', 
                'selModelExtra': 0, 
                'valueMult': 0}
    listKeys = list(dictKeys.keys())
    dictButts = {'buttDivPg': ['Divisão/blocos', ':material/splitscreen:', 'Dividindo o arquivo ', 
                               'Divide o arquivo de acordo com o intervalo de páginas.'], 
                 'buttTxt': ['Texto', ':material/text_ad:', 'Extraindo texto do arquivo ', 
                             'Extrai texto do arquivo e grava o resultado como txt.'],
                 'buttSel': ['Seleção', ':material/filter_alt:', 'Selecionando do arquivo ', 
                             'Cria novo arquivo pdf com as páginas selecionadas.'], 
                 'buttDel': ['Exclusão/páginas', ':material/scan_delete:', 'Deletando do arquivo ', 
                             'Deleta as páginas selecionadas.'], 
                 'buttClear': ['Limpeza', ':material/cleaning_services:', 'Limpando os campos da tela.', 
                               'Limpa os campos da tela, exceto o arquivo escolhido.'], 
                 'buttUrls': ['URLs', ':material/link:', 'Extraindo links/URLs do arquivo ', 
                              'Pesquisa as URLs existentes no arquivo.'], 
                 'buttImgs': ['Imagens', ':material/image:', 'Extraindo imagens do arquivo', 
                              'Extrai imagens do arquivo do arquivo e grava-as individualmente.'], 
                 'buttSize': ['Divisão/tamanho', ':material/view_comfy:', 'Dividindo por tamanho o arquivo ', 
                              'Divide o arquivo de acordo com o tamanho escolhido.'], 
                 'buttMark': ['Marcação', ':material/approval:', 'Marcando o rodapé do arquivo ', 
                              'Insere marca de água no rodapé do arquivo.'], 
                 'buttInfo': ['Informações', ':material/info:', 'Coligindo informações sobre o arquivo inteiro.', 
                              'Exibe informações sobre o arquivo inteiro.'], 
                 'buttTxtTab': ['Pdf/planilha', ':material/transform:', 'Abrindo janela com formatos de tabela ', 
                                'Converte em formato de tabela para as páginas selecionadas.'], 
                 'buttToWord': ['Pdf/documento', ':material/convert_to_text:', 'Convertendo em Word o arquivo ', 
                                'Converte em formato docx as páginas selecionadas do arquivo.'], 
                 'buttToImg': ['Pdf/imagem', ':material/modeling:', 'Convertendo em imagem (png) o arquivo ', 
                               'Converte em formato jpg as páginas selecionadas.'], 
                 'buttToPower': ['Pdf/slide', ':material/cycle:', 'Convertendo em slide do PowerPoint o arquivo ', 
                                 'Converte em slide do PowerPoint as páginas selecionadas.'], 
                 'buttQrcode': ['QR Code', ':material/qr_code_2:', 'Inserindo QR Code no canto inferior direito do arquivo ', 
                                'Insere QR Code no rodapé das páginas selecionadas.'], 
                 'buttPgs': ['Páginas', ':material/view_list:', 'Exibindo opções de seleção de páginas do arquivo ', 
                             'Exibe opções de seleção de páginas.'],
                 'buttToPerson': ['QR Code', ':material/person_edit:', 'Abrindo campos a preencher para inserção do QR Code', 
                                  'Abre opções para preenchimento do QR Code.'], 
                 'buttRemImage': ['Exclusão/imagens', ':material/folder_off:', 'Removendo todas as imagens do arquivo', 
                                  'Remove todas as imagens das páginas selecionadas.'], 
                 'buttRemWords': ['Exclusão/texto', ':material/clear_all:', 'Removendo todas as ocorrências do texto', 
                                  'Remove o texto das páginas selecionadas.'], 
                 'buttOptWords': ['Texto', ':material/text_ad:', 'Abrindo tela para inserção de senha ou de texto a ser substituído', 
                                  'Abre tela para digitar senha ou texto a ser apagado.'], 
                 'buttCodify': ['Bloqueio', ':material/lock:', 'Bloqueando o arquivo', 
                                'Cria senha de bloqueio para o arquivo criado com as´páginas selecionadas.'], 
                 'buttDeCodify': ['Desbloqueio', ':material/lock_open_right:', 'Desbloqueando o arquivo', 
                                  'Desbloqueia todas as páginas do arquivo.'], 
                 'buttNoMark': ['Exclusão/marcas', ':material/variable_remove:', 'Removendo as marcas de água do arquivo', 
                                "Cria arquivo com as´páginas selecionadas e sem marca d'água."], 
                 'buttTypeImgs': ['Planilhas', ':material/format_list_bulleted:', 'Abrindo janela para escolha de opções de imagem.', 
                                  'Abre janela para escolha de tipos de imagem.'], 
                 'buttOptImgs': ['Imagens', ':material/checklist:', 'Abrindo janela para seleção de opções de imagem ', 
                                 'Abre janela com formato de imagem para as páginas selecionadas.'], 
                 'buttOptDocs': ['Documento', ':material/table:', 'Abrindo janela com opções de documento ', 
                                 'Abre janela com formato de documento para as páginas selecioandas.'], 
                 'buttOptSlides': ['Slides', ':material/event_list:', 'Abrindo janela com opções de slide ', 
                                   'Abre janela com formato de slide para as páginas selecionadas.'], 
                 'buttRotate': ['Rotação', ':material/rotate_auto:', 'Abrindo janela com opções de rotação ', 
                                 'Abre janela com opções de páginas e rotação.'], 
                 'buttTxtHtml': ['HTML', ':material/code_blocks:', 'Extraindo texto do arquivo ', 
                                 'Extrai conteúdo HTML e grava o resultado como txt.'], 
                 'buttFinal': ['Rodapé', icons[5][0], 'Indo para o final da página ', 
                               'Rola até o final da página do aplicativo.'], 
                 'buttIni': ['Topo', icons[0][0], 'Indo para o topo da página ', 
                             'Rola até o topo da página do aplicativo.'], 
                 'buttBasic': ['Apoio', icons[1][0], 'Indo para o bloco "Telas de apoio"', 
                               'Rola até o bloco "Telas de apoio" do aplicativo.'], 
                 'buttExtract': ['Extração', icons[2][0], 'Indo para o bloco "Extração de objetos"', 
                                 'Rola até o bloco "Extração de objetos" do aplicativo.'], 
                 'buttOther': ['Operações', icons[3][0], 'Indo para o bloco "Outras operações"', 
                               'Rola até o bloco "Outras operações" do aplicativo.'], 
                 'buttOyher': ['Conversão', icons[4][0], 'Indo para o bloco "Conversão de formatos"', 
                               'Rola até o bloco "Conversão de formatos" do aplicativo.']}
    keysButts = list(dictButts.keys())
    countPg = []
    namesTeste = []
    dirBin = r'C:\Users\ACER\Documents\bin'
    valuesReserve = ['xxxxxxxx xxxxxxx', '(xx)xxxxx-xxxx', 'xxxxxxxx@xxxx.xxx.xx']
    qrCodeKeys = ['one', 'two', 'three']
    keyWords = ['', '', '']
    keyTables = []
    keyDocs = []
    keyImgs = []
    keySlides = []
    for key in qrCodeKeys:
        if key not in st.session_state:
            st.session_state[key] = ''  
    if keyWords not in st.session_state:
        st.session_state[keyWords] = ['', '', '']
    if keyDocs not in st.session_state:
        st.session_state[keyDocs] = []
    if keyImgs not in st.session_state:
        st.session_state[keyImgs] = []
    if keySlides not in st.session_state:
        st.session_state[keySlides] = []
    st.cache_data.clear() 
    iniFinally(0)
    with open('configuration.css') as f:
        css = f.read()
    st.markdown(f'<style>{css}</style>', unsafe_allow_html=True) 
    main()















